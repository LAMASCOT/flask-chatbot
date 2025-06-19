import os
import re
import base64
import requests
import pdfplumber
import openpyxl
from flask import Flask, request, jsonify, send_from_directory, render_template, session
from flask_session import Session
from werkzeug.utils import secure_filename

from langchain.docstore.document import Document
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import UnstructuredHTMLLoader
from langchain_community.chat_models import ChatOpenAI
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings

from docx import Document as DocxDocument
from pptx import Presentation

# --- Chargement des variables d’environnement depuis l'environnement Render (PAS .env ici) ---
USERNAME = os.environ.get("CONFLUENCE_USERNAME")
API_TOKEN = os.environ.get("CONFLUENCE_API_TOKEN")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")

CONFLUENCE_BASE_URL = "https://smartaps.atlassian.net/wiki"
SPACE_KEY = "SAS"
DOCS_DIR = "./docs"
INDEX_DIR = "./faiss_index"
UPLOAD_FOLDER = "./uploads"
ALLOWED_EXTENSIONS = {"pdf", "txt", "docx", "pptx", "xlsx", "png", "jpg", "jpeg", "gif"}

# --- Vérification des variables d'environnement ---
if not USERNAME or not API_TOKEN or not OPENAI_API_KEY:
    raise EnvironmentError("CONFLUENCE_USERNAME, CONFLUENCE_API_TOKEN ou OPENAI_API_KEY manquants.")

os.environ["OPENAI_API_KEY"] = OPENAI_API_KEY
print(f"[DEBUG] Clé OpenAI chargée : {OPENAI_API_KEY[:5]}...")

headers = {
    "Authorization": "Basic " + base64.b64encode(f"{USERNAME}:{API_TOKEN}".encode()).decode(),
    "Content-Type": "application/json"
}

def get_all_pages(space_key):
    pages, start, limit = [], 0, 100
    print("🔄 Récupération des pages depuis Confluence...")
    while True:
        url = f"{CONFLUENCE_BASE_URL}/rest/api/content?spaceKey={space_key}&limit={limit}&start={start}&expand=body.storage"
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        data = response.json()
        for result in data.get("results", []):
            title = result.get("title", "Sans titre")
            content = result.get("body", {}).get("storage", {}).get("value", "")
            if content:
                pages.append({"id": result["id"], "title": title, "content": content})
        if "_links" not in data or "next" not in data["_links"]:
            break
        start += limit
    return pages

def save_pages_as_html(pages):
    os.makedirs(DOCS_DIR, exist_ok=True)
    print("💾 Sauvegarde des pages HTML...")
    for page in pages:
        safe_title = re.sub(r'[<>:"/\\|?*]', '_', page['title'])
        filepath = os.path.join(DOCS_DIR, f"{safe_title}.html")
        try:
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(page['content'])
        except Exception as e:
            print(f"[Erreur] Impossible de sauvegarder {filepath} : {e}")

def create_or_load_faiss_index():
    embeddings = OpenAIEmbeddings(api_key=OPENAI_API_KEY)
    if os.path.exists(INDEX_DIR):
        print("📦 Chargement de l'index FAISS existant...")
        return FAISS.load_local(INDEX_DIR, embeddings, allow_dangerous_deserialization=True)

    print("📤 Création d’un nouvel index FAISS...")
    valid_files = [os.path.join(DOCS_DIR, f) for f in os.listdir(DOCS_DIR) if f.endswith(".html") and os.path.getsize(os.path.join(DOCS_DIR, f)) > 0]
    if not valid_files:
        raise FileNotFoundError("Aucun fichier HTML valide trouvé dans le dossier docs.")

    documents = []
    for filepath in valid_files:
        try:
            loader = UnstructuredHTMLLoader(filepath)
            documents.append(loader.load()[0])
        except Exception as e:
            print(f"[Erreur] Ignoré {filepath} : {e}")

    documents = [doc for doc in documents if doc.page_content.strip()]
    if not documents:
        raise ValueError("Aucun document HTML exploitable après le chargement.")

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = splitter.split_documents(documents)
    db = FAISS.from_documents(docs, embeddings)
    db.save_local(INDEX_DIR)
    return db

app = Flask(__name__)
app.config["SECRET_KEY"] = "votre_cle_secrete_super_securisee"
app.config["SESSION_TYPE"] = "filesystem"
app.config["SESSION_FILE_DIR"] = "./flask_session"
app.config["SESSION_PERMANENT"] = False
Session(app)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

vectorstore = None
qa_chain = None

try:
    if not os.path.exists(INDEX_DIR):
        pages = get_all_pages(SPACE_KEY)
        print(f"✅ {len(pages)} pages récupérées.")
        save_pages_as_html(pages)
    vectorstore = create_or_load_faiss_index()
    qa_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(model_name="gpt-4", temperature=0),
        retriever=vectorstore.as_retriever()
    )
    print("✅ Chaîne QA initialisée.")
except Exception as e:
    print(f"[Erreur d’initialisation] : {e}")

def extract_text_from_pdf(filepath):
    text = ""
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
            else:
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        text += " | ".join(cell.strip() if cell else "" for cell in row) + "\n"
    return text

def extract_text_from_docx(filepath):
    doc = DocxDocument(filepath)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            full_text.append(" | ".join(row_text))
    return "\n".join(full_text)

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, data_only=True)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text.append(" | ".join(row_text))
    return "\n".join(text)

@app.route("/", methods=["GET"])
def home():
    return render_template("index.html")

@app.route('/favicon.ico')
def favicon():
    return send_from_directory(os.path.join(app.root_path, 'static'), 'favicon.ico', mimetype='image/vnd.microsoft.icon')

@app.route("/ask", methods=["POST"])
def ask():
    if not qa_chain:
        return jsonify({"error": "La chaîne QA n’est pas initialisée."}), 500
    try:
        data = request.get_json(force=True)
        question = data.get("question")
        if not question:
            return jsonify({"error": "La question est requise."}), 400
        answer = qa_chain.run(question)
        history = session.get("history", [])
        history.append({"question": question, "answer": answer})
        session["history"] = history
        return jsonify({"answer": answer})
    except Exception as e:
        return jsonify({"error": f"Erreur lors du traitement : {str(e)}"}), 500

@app.route('/upload', methods=['POST'])
def upload_file():
    global vectorstore, qa_chain
    if 'file' not in request.files:
        return jsonify({"error": "Aucun fichier envoyé."}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Nom de fichier vide."}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        try:
            if filename.lower().endswith('.pdf'):
                text = extract_text_from_pdf(filepath)
            elif filename.lower().endswith('.docx'):
                text = extract_text_from_docx(filepath)
            elif filename.lower().endswith('.pptx'):
                text = extract_text_from_pptx(filepath)
            elif filename.lower().endswith('.xlsx'):
                text = extract_text_from_xlsx(filepath)
            else:
                return jsonify({"filename": filename, "message": "Fichier sauvegardé (traitement non implémenté)."})
            if not text.strip():
                return jsonify({"error": "Le fichier est vide ou non extractible."}), 400
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
            docs = splitter.split_text(text)
            documents = [Document(page_content=chunk) for chunk in docs]
            embeddings = OpenAIEmbeddings()
            if vectorstore:
                vectorstore.add_documents(documents)
            else:
                vectorstore = FAISS.from_documents(documents, embeddings)
            vectorstore.save_local(INDEX_DIR)
            qa_chain = RetrievalQA.from_chain_type(
                llm=ChatOpenAI(model_name="gpt-4", temperature=0),
                retriever=vectorstore.as_retriever()
            )
            print(f"[INFO] Fichier '{filename}' indexé avec succès.")
            return jsonify({"filename": filename, "message": f"{filename} indexé avec succès."})
        except Exception as e:
            return jsonify({"error": f"Erreur lors de l’extraction : {str(e)}"}), 500
    return jsonify({"error": "Type de fichier non autorisé."}), 400

@app.route("/history", methods=["GET"])
def get_history():
    return jsonify({"history": session.get("history", [])})

@app.route("/history/clear", methods=["POST"])
def clear_history():
    session.pop("history", None)
    return jsonify({"message": "Historique supprimé avec succès."})

if __name__ == "__main__":
    debug_mode = os.environ.get("FLASK_DEBUG", "false").lower() == "true"
    app.run(host="0.0.0.0", port=5000, debug=debug_mode)
